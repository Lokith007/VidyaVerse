import os
from fastapi import FastAPI, UploadFile, File
from fastapi.responses import FileResponse
from openai import OpenAI
from dotenv import load_dotenv
from reader import extract_notes  # your text extraction code
from pptx import Presentation
import json

# Load environment variables
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

app = FastAPI()

# ---------------- PPTX creation ----------------
def create_ppt_from_json(slides_json, output_file="slides.pptx"):
    prs = Presentation()
    for slide in slides_json:
        title = slide.get("title", "Slide")
        content = slide.get("content", "")
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])
        ppt_slide.shapes.title.text = title
        ppt_slide.placeholders[1].text = content
    prs.save(output_file)
    return output_file

# ---------------- Root endpoint ----------------
@app.get("/")
def root():
    return {"message": "Backend is running!"}

# ---------------- Generate slides ----------------
@app.post("/generate-slides/")
async def generate_slides(file: UploadFile = File(...)):
    # Save uploaded file temporarily
    temp_file_path = f"temp_{file.filename}"
    file_content = await file.read()
    with open(temp_file_path, "wb") as f:
        f.write(file_content)

    try:
        # Extract text using reader.py
        notes_text = extract_notes(temp_file_path)

        # Ask OpenAI to generate slides in JSON format
        prompt = f"""
        Convert the following notes into a JSON list of slides.
        Each slide should have "title" and "content".
        Keep content concise and suitable for slides.
        Notes:
        {notes_text}
        """
        response = client.chat.completions.create(
            model="gpt-4o-mini",   # change to gpt-3.5-turbo if quota issues
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1000,
            temperature=0.5
        )

        slides_text = response.choices[0].message.content

        # Try parsing JSON output from GPT
        try:
            slides_json = json.loads(slides_text)
        except:
            slides_json = [{"title": "Slide 1", "content": slides_text}]

        # Create PPTX file
        ppt_file = f"{file.filename.split('.')[0]}_slides.pptx"
        create_ppt_from_json(slides_json, ppt_file)

    finally:
        # Cleanup temp uploaded file
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)

    # Return PPTX for direct download
    return FileResponse(
        ppt_file,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=ppt_file
    )

